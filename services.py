import os
import json
import pandas as pd  
import docx
import pdfplumber
from openai import OpenAI
from pptx import Presentation
from sqlalchemy.orm import Session
from io import BytesIO
import numpy as np
import cv2
from io import BytesIO
import pytesseract

import crud
from models import DocumentKnowledge, Message
import vector_service


client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


def classify_intent(user_text):
    prompt = f"""تصنف رسالة المستخدم لـ 'greeting' أو 'technical' أو 'out_of_scope'.
    - 'greeting': سلام أو ترحيب.
    - 'technical': سؤال عن العقارات، المنصة، الملفات المرفوعة، بنود العقود، أو الإجراءات القانونية العقارية.
    - 'out_of_scope': مواضيع شخصية، عاطفية، سياسية، أو أي شيء خارج العقارات.

    رد بـ JSON: 
    {{
      "intent": "...",
      "reply": "إذا كان out_of_scope، رد باعتذار مهني يوضح أنك خبير عقاري فقط. إذا كان greeting، رد بترحيب مهني."
    }}
    
    الرسالة: {user_text}"""
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        response_format={ "type": "json_object" },
        temperature=0
    )
    return json.loads(response.choices[0].message.content)

#-------------------------------------------------------------------------------------------

def extract_text_from_pdf(file_content: bytes):
    try:
        with pdfplumber.open(BytesIO(file_content)) as pdf:
            return "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
    except Exception as e:
        print(f"Error PDF: {e}")
        return ""
    
#-------------------------------------------------------------------------------------------

def extract_text_general(content: bytes, filename: str):
    filename = filename.lower()
    try:
        if filename.endswith('.pdf'):
            return extract_text_from_pdf(content)
        elif filename.endswith('.docx'):
            doc = docx.Document(BytesIO(content))
            return "\n".join([p.text for p in doc.paragraphs])
        elif filename.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(BytesIO(content))
            return "بيانات الجدول المستخرجة:\n" + df.to_string()
        elif filename.endswith('.pptx'):
            prs = Presentation(BytesIO(content))
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    except Exception as e:
        print(f"Extraction Error for {filename}: {e}")
    return ""

#-------------------------------------------------------------------------------------------

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

async def extract_text_from_image(image_bytes: bytes):
    try:
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        
        text = pytesseract.image_to_string(gray, lang='ara+eng')
        
        if not text.strip():
            return "لم يتم العثور على نص واضح."
            
        return text
    except Exception as e:
        return f"Tesseract Error: {str(e)}"
    
#-------------------------------------------------------------------------------------------

def get_ai_answer(db: Session, question: str, conv):

    # ----------- Vector search (site knowledge)
    manual_text = vector_service.search_vector_db(
        question,
        collection_type="site"
    )

    # ----------- Conversation documents
    db_docs = db.query(DocumentKnowledge)\
        .filter(DocumentKnowledge.conversation_id == conv.id)\
        .all()

    docs_list = []
    for d in db_docs:
        if d.content and len(d.content) > 20 and d.content != "PROCESSING":
            docs_list.append(
                f"--- بداية المستند ({d.file_name}) ---\n"
                f"{d.content}\n"
                f"--- نهاية المستند ---"
            )

    docs_text = "\n\n".join(docs_list)

    # ----------- Chat history
    history = crud.get_chat_history(db, conv.id, limit=5)

    chat_context = [
        {"role": m.role, "content": m.text}
        for m in reversed(history)
        if m.text
    ]

    # ----------- System prompt
    system_prompt = f"""
أنت خبير عقاري ذكي للمنصة الوطنية للعقارات.
لديك وثائق مرفوعة من قبل المستخدم، وقد تحتوي على أخطاء بسيطة بسبب OCR.
مهمتك هي تحليل الوثائق والإجابة على الأسئلة بدقة ومهنية.

[نصوص الوثائق المرفوعة في المحادثة]:
{docs_text if docs_text else "لا توجد وثائق مرفوعة حالياً."}

[دليل المنصة العام]:
{manual_text}

تعليمات:
1. إذا ذُكر "هذا العقد" أو "الصورة"، ارجع أولاً لنصوص الوثائق.
2. صحح الأخطاء الإملائية الشائعة إن وُجدت.
3. إذا بدا المستند عقداً رسمياً، وضّح ذلك للمستخدم.
4. فقط إذا لم تجد أي معلومة مفيدة إطلاقاً، اكتب NOT_FOUND.
""".strip()

    messages = [{"role": "system", "content": system_prompt}]
    messages.extend(chat_context)
    messages.append({"role": "user", "content": question})

    # ----------- Placeholder message (IMPORTANT)
    assistant_msg = Message(
        conversation_id=conv.id,
        role="assistant",
        type="text",
        text=""
    )
    db.add(assistant_msg)
    db.commit()
    db.refresh(assistant_msg)

    full_response = ""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0,
            stream=True
        )

        for chunk in response:
            if chunk.choices and chunk.choices[0].delta.content:
                token = chunk.choices[0].delta.content
                full_response += token
                yield token

        # ----------- Finalize message
        if "NOT_FOUND" in full_response:
            crud.add_unanswered(db, question)
            full_response += "\n\n(ملاحظة: عذراً، لم أجد هذه التفاصيل، يرجى التواصل مع الدعم)."

        assistant_msg.text = full_response
        db.commit()

    except Exception as e:
        print(f"Error in AI Response: {e}")
        assistant_msg.text = "عذراً، حدث خطأ أثناء الاتصال بالذكاء الاصطناعي."
        db.commit()
        yield assistant_msg.text


#-------------------------------------------------------------------------------------------

def generate_chat_title(first_question: str):
    prompt = f"صغ عنواناً جذاباً وقصيراً جداً (3 كلمات) لهذا السؤال: {first_question}"
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=20
        )
        return response.choices[0].message.content.strip()
    except:
        return first_question[:30]
    
#-------------------------------------------------------------------------------------------
    
async def process_file_task(doc_id: int, file_path: str, filename: str, conv_id: int):
    from database import SessionLocal 
    import database 
    db = SessionLocal()
    try:
        with open(file_path, "rb") as f:
            content = f.read()

        if filename.lower().endswith(('.jpg', '.jpeg', '.png')):
            text = await extract_text_from_image(content)
        else:
            text = extract_text_general(content, filename)
        
        if text:
            doc = db.query(database.DocumentKnowledge).filter(database.DocumentKnowledge.id == doc_id).first()
            if doc:
                doc.content = text
                db.commit()
    except Exception as e:
        print(f"Error in background task: {e}")
    finally:
        db.close()